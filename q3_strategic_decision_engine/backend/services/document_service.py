"""
Document Processing Service
Handles file uploads, text extraction, and document processing for the strategic planning platform
"""

import asyncio
import os
import uuid
from typing import List, Dict, Any, Optional, Tuple
import logging
from dataclasses import dataclass
from pathlib import Path
import mimetypes
import hashlib
import time

# Document processing imports
import pypdf
from docx import Document as DocxDocument
from pptx import Presentation
import pandas as pd
from langchain_core.documents import Document

from ..core.config import settings
from ..core.database import DatabaseManager
from ..core.logging_config import get_logger


@dataclass
class ProcessedDocument:
    """Processed document with metadata"""
    document_id: str
    filename: str
    file_path: str
    content: str
    metadata: Dict[str, Any]
    chunks: List[Document]
    processing_time: float
    summary: Optional[str] = None
    key_topics: Optional[List[str]] = None


@dataclass
class DocumentStats:
    """Document statistics"""
    total_documents: int
    total_chunks: int
    total_size: int
    processing_time: float
    file_types: Dict[str, int]


class DocumentService:
    """Service for processing documents and managing file uploads"""
    
    def __init__(self, vector_store_service=None, llm_service=None):
        self.vector_store_service = vector_store_service
        self.llm_service = llm_service
        self.logger = get_logger('document')
        self.upload_dir = Path(settings.UPLOAD_DIR)
        self.upload_dir.mkdir(parents=True, exist_ok=True)
    
    async def upload_file(self, file_content: bytes, filename: str, content_type: str = None) -> Dict[str, Any]:
        """Upload and save file"""
        try:
            # Validate file
            validation_result = self._validate_file(filename, len(file_content), content_type)
            if not validation_result['valid']:
                return {
                    'success': False,
                    'error': validation_result['error'],
                    'document_id': None
                }
            
            # Generate unique filename
            file_extension = Path(filename).suffix.lower()
            document_id = str(uuid.uuid4())
            unique_filename = f"{document_id}{file_extension}"
            file_path = self.upload_dir / unique_filename
            
            # Save file
            with open(file_path, 'wb') as f:
                f.write(file_content)
            
            # Create database record
            with DatabaseManager() as db:
                document = db.create_document(
                    filename=unique_filename,
                    original_filename=filename,
                    file_path=str(file_path),
                    file_size=len(file_content),
                    file_type=file_extension,
                    mime_type=content_type or mimetypes.guess_type(filename)[0],
                    metadata={
                        'upload_timestamp': time.time(),
                        'file_hash': hashlib.md5(file_content).hexdigest()
                    }
                )
                document_id = document.id
            
            self.logger.info(f"File uploaded successfully: {filename} -> {unique_filename}")
            
            return {
                'success': True,
                'document_id': document_id,
                'filename': unique_filename,
                'original_filename': filename,
                'file_size': len(file_content),
                'file_type': file_extension,
                'error': None
            }
            
        except Exception as e:
            self.logger.error(f"File upload failed: {str(e)}")
            return {
                'success': False,
                'error': f"Upload failed: {str(e)}",
                'document_id': None
            }
    
    def _validate_file(self, filename: str, file_size: int, content_type: str = None) -> Dict[str, Any]:
        """Validate uploaded file"""
        try:
            # Check file extension
            file_extension = Path(filename).suffix.lower().lstrip('.')
            if file_extension not in settings.ALLOWED_EXTENSIONS:
                return {
                    'valid': False,
                    'error': f"File type '{file_extension}' not allowed. Allowed types: {', '.join(settings.ALLOWED_EXTENSIONS)}"
                }
            
            # Check file size
            if file_size > settings.MAX_FILE_SIZE:
                return {
                    'valid': False,
                    'error': f"File size ({file_size} bytes) exceeds maximum allowed size ({settings.MAX_FILE_SIZE} bytes)"
                }
            
            # Check if file is empty
            if file_size == 0:
                return {
                    'valid': False,
                    'error': "File is empty"
                }
            
            return {'valid': True, 'error': None}
            
        except Exception as e:
            return {
                'valid': False,
                'error': f"Validation error: {str(e)}"
            }
    
    async def process_document(self, document_id: int, generate_summary: bool = True) -> ProcessedDocument:
        """Process document and extract content"""
        start_time = time.time()
        
        try:
            # Get document from database
            with DatabaseManager() as db:
                document = db.get_document(document_id)
                if not document:
                    raise ValueError(f"Document with ID {document_id} not found")
            
            # Extract content based on file type
            content, metadata = await self._extract_content(document.file_path, document.file_type)
            
            if not content.strip():
                raise ValueError("No content could be extracted from the document")
            
            # Create LangChain document
            langchain_doc = Document(
                page_content=content,
                metadata={
                    'source': document.original_filename,
                    'document_id': str(document_id),
                    'file_type': document.file_type,
                    'file_size': document.file_size,
                    **metadata
                }
            )
            
            # Add to vector store
            chunks_count = 0
            if self.vector_store_service:
                chunks_count = await self.vector_store_service.add_documents(
                    [langchain_doc],
                    document_id=str(document_id)
                )
            
            # Generate summary and key topics if requested
            summary = None
            key_topics = None
            
            if generate_summary and self.llm_service:
                summary = await self._generate_summary(content)
                key_topics = await self._extract_key_topics(content)
            
            # Update database
            with DatabaseManager() as db:
                db.update_document_processed(
                    document_id=document_id,
                    chunk_count=chunks_count,
                    summary=summary,
                    key_topics=key_topics
                )
            
            processing_time = time.time() - start_time
            
            self.logger.info(f"Document processed successfully: {document.original_filename} ({chunks_count} chunks)")
            
            return ProcessedDocument(
                document_id=str(document_id),
                filename=document.original_filename,
                file_path=document.file_path,
                content=content,
                metadata=metadata,
                chunks=[langchain_doc],
                processing_time=processing_time,
                summary=summary,
                key_topics=key_topics
            )
            
        except Exception as e:
            self.logger.error(f"Document processing failed: {str(e)}")
            raise
    
    async def _extract_content(self, file_path: str, file_type: str) -> Tuple[str, Dict[str, Any]]:
        """Extract content from file based on type"""
        try:
            content = ""
            metadata = {}
            
            if file_type == '.pdf':
                content, metadata = await self._extract_pdf_content(file_path)
            elif file_type == '.docx':
                content, metadata = await self._extract_docx_content(file_path)
            elif file_type == '.pptx':
                content, metadata = await self._extract_pptx_content(file_path)
            elif file_type == '.xlsx':
                content, metadata = await self._extract_xlsx_content(file_path)
            elif file_type in ['.txt', '.csv']:
                content, metadata = await self._extract_text_content(file_path)
            else:
                raise ValueError(f"Unsupported file type: {file_type}")
            
            return content, metadata
            
        except Exception as e:
            self.logger.error(f"Content extraction failed for {file_path}: {str(e)}")
            raise
    
    async def _extract_pdf_content(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Extract content from PDF file"""
        try:
            content = ""
            metadata = {'pages': 0, 'has_images': False}
            
            with open(file_path, 'rb') as file:
                pdf_reader = pypdf.PdfReader(file)
                metadata['pages'] = len(pdf_reader.pages)
                
                for page_num, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text()
                    if page_text.strip():
                        content += f"\n--- Page {page_num + 1} ---\n"
                        content += page_text
                    
                    # Check for images
                    if '/XObject' in page.get('/Resources', {}):
                        metadata['has_images'] = True
                
                # Extract metadata
                if pdf_reader.metadata:
                    metadata.update({
                        'title': pdf_reader.metadata.get('/Title', ''),
                        'author': pdf_reader.metadata.get('/Author', ''),
                        'subject': pdf_reader.metadata.get('/Subject', ''),
                        'creator': pdf_reader.metadata.get('/Creator', ''),
                        'producer': pdf_reader.metadata.get('/Producer', ''),
                        'creation_date': pdf_reader.metadata.get('/CreationDate', ''),
                        'modification_date': pdf_reader.metadata.get('/ModDate', '')
                    })
            
            return content.strip(), metadata
            
        except Exception as e:
            raise Exception(f"PDF extraction failed: {str(e)}")
    
    async def _extract_docx_content(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Extract content from DOCX file"""
        try:
            doc = DocxDocument(file_path)
            content = ""
            metadata = {'paragraphs': 0, 'tables': 0, 'images': 0}
            
            # Extract paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    content += paragraph.text + "\n"
                    metadata['paragraphs'] += 1
            
            # Extract tables
            for table in doc.tables:
                metadata['tables'] += 1
                content += "\n--- Table ---\n"
                for row in table.rows:
                    row_text = " | ".join(cell.text for cell in row.cells)
                    content += row_text + "\n"
            
            # Count images
            from docx.document import Document as DocxDoc
            metadata['images'] = len(doc.inline_shapes)
            
            # Extract core properties
            if doc.core_properties:
                metadata.update({
                    'title': doc.core_properties.title or '',
                    'author': doc.core_properties.author or '',
                    'subject': doc.core_properties.subject or '',
                    'created': str(doc.core_properties.created) if doc.core_properties.created else '',
                    'modified': str(doc.core_properties.modified) if doc.core_properties.modified else ''
                })
            
            return content.strip(), metadata
            
        except Exception as e:
            raise Exception(f"DOCX extraction failed: {str(e)}")
    
    async def _extract_pptx_content(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Extract content from PPTX file"""
        try:
            prs = Presentation(file_path)
            content = ""
            metadata = {'slides': 0, 'text_shapes': 0, 'images': 0}
            
            for slide_num, slide in enumerate(prs.slides):
                metadata['slides'] += 1
                content += f"\n--- Slide {slide_num + 1} ---\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content += shape.text + "\n"
                        metadata['text_shapes'] += 1
                    
                    # Count images
                    if shape.shape_type == 13:  # Picture
                        metadata['images'] += 1
            
            # Extract core properties
            if prs.core_properties:
                metadata.update({
                    'title': prs.core_properties.title or '',
                    'author': prs.core_properties.author or '',
                    'subject': prs.core_properties.subject or '',
                    'created': str(prs.core_properties.created) if prs.core_properties.created else '',
                    'modified': str(prs.core_properties.modified) if prs.core_properties.modified else ''
                })
            
            return content.strip(), metadata
            
        except Exception as e:
            raise Exception(f"PPTX extraction failed: {str(e)}")
    
    async def _extract_xlsx_content(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Extract content from XLSX file"""
        try:
            # Read all sheets
            xlsx_file = pd.ExcelFile(file_path)
            content = ""
            metadata = {'sheets': len(xlsx_file.sheet_names), 'total_rows': 0, 'total_columns': 0}
            
            for sheet_name in xlsx_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                content += f"\n--- Sheet: {sheet_name} ---\n"
                
                # Convert to string representation
                content += df.to_string(index=False)
                content += "\n"
                
                metadata['total_rows'] += len(df)
                metadata['total_columns'] = max(metadata['total_columns'], len(df.columns))
            
            return content.strip(), metadata
            
        except Exception as e:
            raise Exception(f"XLSX extraction failed: {str(e)}")
    
    async def _extract_text_content(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Extract content from text files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
            
            metadata = {
                'lines': len(content.split('\n')),
                'characters': len(content),
                'words': len(content.split())
            }
            
            return content, metadata
            
        except UnicodeDecodeError:
            # Try different encodings
            for encoding in ['latin-1', 'cp1252', 'iso-8859-1']:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        content = file.read()
                    
                    metadata = {
                        'lines': len(content.split('\n')),
                        'characters': len(content),
                        'words': len(content.split()),
                        'encoding': encoding
                    }
                    
                    return content, metadata
                except UnicodeDecodeError:
                    continue
            
            raise Exception("Unable to decode text file with supported encodings")
        except Exception as e:
            raise Exception(f"Text extraction failed: {str(e)}")
    
    async def _generate_summary(self, content: str) -> Optional[str]:
        """Generate document summary using LLM"""
        try:
            if not self.llm_service:
                return None
            
            # Limit content length for summary generation
            max_length = 4000
            if len(content) > max_length:
                content = content[:max_length] + "..."
            
            prompt = f"""Please provide a concise summary of the following document content. 
            Focus on the main points, key information, and overall purpose of the document.
            
            Document content:
            {content}
            
            Summary:"""
            
            response = await self.llm_service.generate_response(
                prompt=prompt,
                task=self.llm_service.LLMTask.DOCUMENT_ANALYSIS,
                max_tokens=200
            )
            
            return response.content if response else None
            
        except Exception as e:
            self.logger.error(f"Summary generation failed: {str(e)}")
            return None
    
    async def _extract_key_topics(self, content: str) -> Optional[List[str]]:
        """Extract key topics from document using LLM"""
        try:
            if not self.llm_service:
                return None
            
            # Limit content length
            max_length = 4000
            if len(content) > max_length:
                content = content[:max_length] + "..."
            
            prompt = f"""Analyze the following document content and extract the main topics and themes.
            Return a list of 5-10 key topics or keywords that best represent the content.
            Format as a simple comma-separated list.
            
            Document content:
            {content}
            
            Key topics:"""
            
            response = await self.llm_service.generate_response(
                prompt=prompt,
                task=self.llm_service.LLMTask.DOCUMENT_ANALYSIS,
                max_tokens=100
            )
            
            if response and response.content:
                # Parse topics from response
                topics = [topic.strip() for topic in response.content.split(',')]
                return topics[:10]  # Limit to 10 topics
            
            return None
            
        except Exception as e:
            self.logger.error(f"Key topics extraction failed: {str(e)}")
            return None
    
    async def get_document_info(self, document_id: int) -> Optional[Dict[str, Any]]:
        """Get document information"""
        try:
            with DatabaseManager() as db:
                document = db.get_document(document_id)
                if not document:
                    return None
                
                return {
                    'id': document.id,
                    'filename': document.original_filename,
                    'file_size': document.file_size,
                    'file_type': document.file_type,
                    'upload_date': document.upload_date.isoformat(),
                    'processed': document.processed,
                    'processed_date': document.processed_date.isoformat() if document.processed_date else None,
                    'chunk_count': document.chunk_count,
                    'summary': document.summary,
                    'key_topics': document.key_topics,
                    'metadata': document.metadata
                }
                
        except Exception as e:
            self.logger.error(f"Failed to get document info: {str(e)}")
            return None
    
    async def list_documents(self, skip: int = 0, limit: int = 100) -> List[Dict[str, Any]]:
        """List all documents"""
        try:
            with DatabaseManager() as db:
                documents = db.get_documents(skip=skip, limit=limit)
                
                return [
                    {
                        'id': doc.id,
                        'filename': doc.original_filename,
                        'file_size': doc.file_size,
                        'file_type': doc.file_type,
                        'upload_date': doc.upload_date.isoformat(),
                        'processed': doc.processed,
                        'chunk_count': doc.chunk_count,
                        'summary': doc.summary[:200] + "..." if doc.summary and len(doc.summary) > 200 else doc.summary
                    }
                    for doc in documents
                ]
                
        except Exception as e:
            self.logger.error(f"Failed to list documents: {str(e)}")
            return []
    
    async def delete_document(self, document_id: int) -> bool:
        """Delete document and its chunks"""
        try:
            # Get document info
            with DatabaseManager() as db:
                document = db.get_document(document_id)
                if not document:
                    return False
            
            # Delete from vector store
            if self.vector_store_service:
                await self.vector_store_service.delete_document(str(document_id))
            
            # Delete file
            try:
                if os.path.exists(document.file_path):
                    os.remove(document.file_path)
            except Exception as e:
                self.logger.warning(f"Failed to delete file {document.file_path}: {str(e)}")
            
            # Mark as inactive in database
            with DatabaseManager() as db:
                document = db.get_document(document_id)
                if document:
                    document.is_active = False
                    db.db.commit()
            
            self.logger.info(f"Document {document_id} deleted successfully")
            return True
            
        except Exception as e:
            self.logger.error(f"Failed to delete document: {str(e)}")
            return False
    
    async def get_stats(self) -> DocumentStats:
        """Get document processing statistics"""
        try:
            with DatabaseManager() as db:
                documents = db.get_documents(limit=1000)  # Get a reasonable sample
                
                total_documents = len(documents)
                total_chunks = sum(doc.chunk_count for doc in documents)
                total_size = sum(doc.file_size for doc in documents)
                total_processing_time = 0  # This would need to be tracked separately
                
                file_types = {}
                for doc in documents:
                    file_type = doc.file_type
                    file_types[file_type] = file_types.get(file_type, 0) + 1
                
                return DocumentStats(
                    total_documents=total_documents,
                    total_chunks=total_chunks,
                    total_size=total_size,
                    processing_time=total_processing_time,
                    file_types=file_types
                )
                
        except Exception as e:
            self.logger.error(f"Failed to get document stats: {str(e)}")
            return DocumentStats(0, 0, 0, 0.0, {}) 